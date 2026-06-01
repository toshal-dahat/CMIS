import base64
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _collect_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return []

    blocks = []
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        runs = [run.text for run in paragraph.runs if run.text]
        paragraph_text = "".join(runs).strip()
        if paragraph_text:
            blocks.append(paragraph_text)
    return blocks


def _collect_table(shape):
    if not getattr(shape, "has_table", False):
        return []

    rows = []
    for row in shape.table.rows:
        cells = []
        for cell in row.cells:
            cell_text = " ".join(part.strip() for part in cell.text.splitlines() if part.strip()).strip()
            cells.append(cell_text)
        if any(cells):
            rows.append(" | ".join(cells))
    return rows


def _collect_chart(shape):
    if not getattr(shape, "has_chart", False):
        return []

    chart = shape.chart
    categories = []
    try:
        first_plot = chart.plots[0]
        categories = [str(category.label) for category in first_plot.categories]
    except Exception:
        categories = []

    summaries = []
    for series in chart.series:
        try:
            values = list(series.values)
        except Exception:
            values = []

        series_name = str(series.name or "Series")
        if categories and len(categories) == len(values):
            pairs = ", ".join(f"{categories[index]}: {values[index]}" for index in range(len(values)))
            summaries.append(f"{series_name} -> {pairs}")
        elif values:
            summaries.append(f"{series_name} -> {', '.join(str(value) for value in values)}")
        else:
            summaries.append(series_name)
    return summaries


def _collect_picture(shape, slide_number):
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None

    image = shape.image
    blob = image.blob
    return {
        "slideNumber": slide_number,
        "name": getattr(shape, "name", f"slide-{slide_number}-image"),
        "contentType": image.content_type,
        "extension": image.ext,
        "byteSize": len(blob),
        "sha1": getattr(image, "sha1", None),
        "base64Data": base64.b64encode(blob).decode("ascii"),
    }


def _collect_notes(slide):
    if not getattr(slide, "has_notes_slide", False):
        return ""

    try:
        notes_frame = slide.notes_slide.notes_text_frame
    except Exception:
        return ""

    if not notes_frame:
        return ""

    return notes_frame.text.strip()


def _find_title(slide):
    title_shape = slide.shapes.title
    if title_shape and getattr(title_shape, "has_text_frame", False):
        title_text = title_shape.text.strip()
        if title_text:
            return title_text
    return f"Slide {slide.slide_id}"


def extract_presentation(pptx_path, max_images_returned=6, max_image_bytes=400 * 1024):
    presentation = Presentation(str(Path(pptx_path)))
    slides_payload = []
    selected_images = []
    seen_image_hashes = set()

    for index, slide in enumerate(presentation.slides, start=1):
        text_blocks = []
        table_rows = []
        chart_summaries = []
        slide_images = []

        for shape in _iter_shapes(slide.shapes):
            text_blocks.extend(_collect_text(shape))
            table_rows.extend(_collect_table(shape))
            chart_summaries.extend(_collect_chart(shape))

            picture = _collect_picture(shape, index)
            if picture:
                slide_images.append(
                    {
                        "name": picture["name"],
                        "contentType": picture["contentType"],
                        "extension": picture["extension"],
                        "byteSize": picture["byteSize"],
                    }
                )
                picture_hash = picture.get("sha1") or f"{picture['slideNumber']}:{picture['name']}"
                if (
                    len(selected_images) < max_images_returned
                    and picture["byteSize"] <= max_image_bytes
                    and picture_hash not in seen_image_hashes
                ):
                    seen_image_hashes.add(picture_hash)
                    selected_images.append(picture)

        notes = _collect_notes(slide)
        slides_payload.append(
            {
                "slideNumber": index,
                "slideId": slide.slide_id,
                "title": _find_title(slide),
                "textBlocks": text_blocks,
                "tableRows": table_rows,
                "chartSummaries": chart_summaries,
                "notes": notes,
                "images": slide_images,
            }
        )

    all_text_fragments = []
    for slide in slides_payload:
        all_text_fragments.extend(slide["textBlocks"])
        all_text_fragments.extend(slide["tableRows"])
        all_text_fragments.extend(slide["chartSummaries"])
        if slide["notes"]:
            all_text_fragments.append(slide["notes"])

    return {
        "slideCount": len(slides_payload),
        "slides": slides_payload,
        "images": selected_images,
        "allText": "\n".join(fragment for fragment in all_text_fragments if fragment).strip(),
    }
